# -*- coding: utf-8 -*-
"""
decklib.py — biblioteca do deck, ESTILO VIVO/VIBRANTE (uma cor por bloco, cartões,
fontes grandes). Gera .pptx + PDF companheiro (imagem do slide + roteiro + tempo).
"""
import os, subprocess, glob, html
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FONT = "Liberation Sans"
MONO = "DejaVu Sans Mono"

# ---------- PALETA VIBRANTE ----------
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x16, 0x1E, 0x33)   # texto principal (quase preto azulado)
MUTED  = RGBColor(0x5A, 0x66, 0x7E)   # texto secundário
PAPER  = RGBColor(0xFB, 0xFC, 0xFE)   # fundo (branco levemente frio)
LINE   = RGBColor(0xE2, 0xE7, 0xF0)
CODEBG = RGBColor(0x12, 0x1A, 0x30)
CODEFG = RGBColor(0xEA, 0xEF, 0xFA)
GOOD   = RGBColor(0x0F, 0x9D, 0x58)
WARN   = RGBColor(0xE0, 0x6A, 0x00)

# cor por família: (forte, médio, tinta clara p/ cartão)
FAM = {
    "indigo":  (RGBColor(0x4F,0x46,0xE5), RGBColor(0x7A,0x74,0xF0), RGBColor(0xEC,0xEB,0xFD)),
    "sky":     (RGBColor(0x02,0x83,0xC9), RGBColor(0x38,0xA8,0xE0), RGBColor(0xE3,0xF3,0xFC)),
    "violet":  (RGBColor(0x7C,0x3A,0xED), RGBColor(0x9B,0x66,0xF0), RGBColor(0xF2,0xEA,0xFE)),
    "teal":    (RGBColor(0x0D,0x94,0x88), RGBColor(0x2C,0xB1,0xA6), RGBColor(0xE1,0xF5,0xF3)),
    "emerald": (RGBColor(0x05,0x96,0x69), RGBColor(0x2E,0xB0,0x85), RGBColor(0xE2,0xF6,0xEE)),
    "amber":   (RGBColor(0xD9,0x77,0x06), RGBColor(0xEA,0x99,0x2E), RGBColor(0xFC,0xF1,0xDF)),
    "rose":    (RGBColor(0xE1,0x1D,0x48), RGBColor(0xEC,0x54,0x76), RGBColor(0xFC,0xE6,0xEC)),
    "slate":   (RGBColor(0x33,0x41,0x66), RGBColor(0x5A,0x66,0x82), RGBColor(0xEA,0xED,0xF4)),
}
# bloco -> família de cor
BLOCK_FAM = {0:"indigo",1:"violet",2:"sky",3:"teal",4:"emerald",5:"amber",
             6:"rose",7:"indigo",8:"teal",9:"indigo",10:"emerald",11:"violet",12:"slate"}

def fam(block):
    return FAM[BLOCK_FAM.get(block, "indigo")]

SW, SH = 13.333, 7.5


def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.06):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if rounded:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def _dot(slide, x, y, d, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    return sp

def _txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.06, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        try: p.line_spacing = line_spacing
        except Exception: pass
        for (t, st) in para:
            r = p.add_run(); r.text = t; f = r.font
            f.name = st.get("font", FONT); f.size = Pt(st.get("size", 18))
            f.bold = st.get("bold", False); f.italic = st.get("italic", False)
            f.color.rgb = st.get("color", INK)
    return tb


def cover(slide, title, subtitle, tagline):
    c1, c2, c3 = FAM["indigo"]
    _rect(slide, 0, 0, SW, SH, fill=PAPER)
    # faixa de cor à esquerda + acentos geométricos
    _rect(slide, 0, 0, 0.55, SH, fill=c1)
    _rect(slide, 0.55, 0, 0.14, SH, fill=c2)
    _dot(slide, SW-2.4, -1.2, 3.6, FAM["sky"][2])
    _dot(slide, SW-1.3, SH-1.6, 2.6, FAM["violet"][2])
    _rect(slide, 1.15, 1.75, 1.4, 0.5, fill=c1, rounded=True, radius=0.5)
    _txt(slide, 1.30, 1.80, 1.2, 0.4, [[("PALESTRA", {"size": 12, "bold": True, "color": WHITE})]],
         anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, 1.15, 2.55, SW-3.0, 2.1,
         [[(title, {"size": 37, "bold": True, "color": c1})]], line_spacing=1.03)
    _rect(slide, 1.18, 4.75, 2.6, 0.09, fill=FAM["amber"][0])
    _txt(slide, 1.15, 5.0, SW-3.2, 1.0,
         [[(subtitle, {"size": 19, "italic": True, "color": INK})]], line_spacing=1.12)
    _txt(slide, 1.15, 6.55, SW-3.2, 0.6,
         [[(tagline, {"size": 13, "color": MUTED})]])


def header(slide, block, kicker, title):
    c1, c2, c3 = fam(block)
    _rect(slide, 0, 0, SW, SH, fill=PAPER)
    _rect(slide, 0, 0, SW, 0.18, fill=c1)                    # faixa superior colorida
    # chip do bloco (uma linha só)
    chip_w = 0.5 + 0.108*len(kicker)
    _rect(slide, 0.6, 0.5, chip_w, 0.44, fill=c1, rounded=True, radius=0.5)
    _txt(slide, 0.72, 0.5, chip_w, 0.44,
         [[(kicker.upper(), {"size": 12.5, "bold": True, "color": WHITE})]],
         anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    # número do slide (chip contornado)
    # título grande
    _txt(slide, 0.6, 1.06, SW-1.2, 0.9, [[(title, {"size": 31, "bold": True, "color": INK})]])
    _rect(slide, 0.62, 1.78, 2.3, 0.09, fill=c2)             # sublinhado curto e colorido
    return (c1, c2, c3)


def footer(slide, block, n, total, label="Engenharia de IA & Desenvolvimento Orientado a Especificação"):
    c1, c2, c3 = fam(block)
    _rect(slide, 0, SH-0.3, SW, 0.3, fill=c3)
    _rect(slide, 0, SH-0.3, 0.12, 0.3, fill=c1)
    _txt(slide, 0.6, SH-0.29, 9.5, 0.27, [[(label, {"size": 10.5, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, SW-1.7, SH-0.29, 1.1, 0.27, [[("%02d / %02d" % (n, total), {"size": 10.5, "bold": True, "color": c1})]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def cards(slide, block, items, x=0.6, y=2.0, w=SW-1.2, size=18.5, gap=0.14, ch=None, lead_color=None):
    """items: (lead, resto)  OU  (lead, resto, kind) kind∈{good,warn}.
       Cartão full-width: borda colorida à esquerda + tinta clara + lead em negrito colorido."""
    c1, c2, c3 = fam(block)
    n = len(items)
    total_h = SH - 0.55 - y
    if ch is None:
        ch = (total_h - (n-1)*gap) / n
        ch = min(ch, 1.15)
    cy = y
    for it in items:
        lead = it[0]; rest = it[1] if len(it) > 1 else ""
        kind = it[2] if len(it) > 2 else None
        bar = {"good": GOOD, "warn": WARN}.get(kind, c1)
        tint = {"good": RGBColor(0xE6,0xF6,0xEC), "warn": RGBColor(0xFC,0xF1,0xE2)}.get(kind, c3)
        _rect(slide, x, cy, w, ch, fill=tint, rounded=True, radius=0.10)
        _rect(slide, x, cy, 0.12, ch, fill=bar)
        runs = [[(lead, {"size": size, "bold": True, "color": (lead_color or bar)}),
                 ("  " + rest if rest else "", {"size": size, "color": INK})]]
        _txt(slide, x+0.35, cy, w-0.6, ch, runs, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05, space_after=0)
        cy += ch + gap


def bullets(slide, items, x=0.6, y=2.0, w=SW-1.2, h=SH-2.6, size=19, gap=9, lh=1.12, block=0):
    c1, c2, c3 = fam(block)
    runs = []
    for it in items:
        lvl = it[0]; parts = it[1]
        pre = ("●  " if lvl == 0 else "–  ")
        seg = [(pre, {"size": size-(1 if lvl else 0), "bold": True, "color": (c1 if lvl == 0 else c2)})]
        for (t, code) in parts:
            st = {"size": size-(1 if lvl else 0), "color": INK}
            if code == "b": st.update(bold=True, color=c1)
            elif code == "m": st.update(color=MUTED)
            elif code == "g": st.update(bold=True, color=GOOD)
            elif code == "w": st.update(bold=True, color=WARN)
            elif code == "code": st.update(font=MONO, size=size-4, color=c1)
            seg.append((t, st))
        runs.append(seg)
    _txt(slide, x, y, w, h, runs, space_after=gap, line_spacing=lh)


def panel(slide, x, y, w, h, block=0, kind=None):
    c1, c2, c3 = fam(block)
    bar = {"good": GOOD, "warn": WARN}.get(kind, c1)
    tint = {"good": RGBColor(0xE6,0xF6,0xEC), "warn": RGBColor(0xFC,0xF1,0xE2)}.get(kind, c3)
    _rect(slide, x, y, w, h, fill=tint, rounded=True, radius=0.08)
    _rect(slide, x, y, 0.12, h, fill=bar)
    return bar, tint


def callout(slide, block, runs_parts, x=0.6, y=None, w=SW-1.2, h=1.05, size=17, kind=None):
    if y is None: y = SH-1.45
    bar, tint = panel(slide, x, y, w, h, block=block, kind=kind)
    _txt(slide, x+0.35, y, w-0.65, h, [runs_parts], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)


def quote(slide, block, text_parts, x=0.6, y=2.05, w=SW-1.2, h=1.6, size=25):
    c1, c2, c3 = fam(block)
    _rect(slide, x, y, w, h, fill=c3, rounded=True, radius=0.06)
    _rect(slide, x, y, 0.16, h, fill=c1)
    _txt(slide, x+0.5, y, w-0.9, h, [text_parts], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)


def table(slide, block, headers, rows, x=0.6, y=2.0, w=SW-1.2, h=None, fsize=13, header_fs=13):
    c1, c2, c3 = fam(block)
    nr = len(rows)+1; nc = len(headers)
    if h is None: h = min(SH-2.5, 0.46*nr + 0.15)
    gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    for ci, htext in enumerate(headers):
        c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = c1
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for m in ("left","right"): setattr(c, "margin_"+m, Inches(0.08))
        for m in ("top","bottom"): setattr(c, "margin_"+m, Inches(0.03))
        tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r = p.add_run(); r.text = htext; r.font.name = FONT; r.font.size = Pt(header_fs)
        r.font.bold = True; r.font.color.rgb = WHITE
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci); c.fill.solid()
            c.fill.fore_color.rgb = c3 if ri % 2 else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            for m in ("left","right"): setattr(c, "margin_"+m, Inches(0.08))
            for m in ("top","bottom"): setattr(c, "margin_"+m, Inches(0.02))
            tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
            bold = val.startswith("*")
            r = p.add_run(); r.text = val[1:] if bold else val
            r.font.name = FONT; r.font.size = Pt(fsize); r.font.bold = bold
            r.font.color.rgb = c1 if bold else INK
    return tbl


def code_block(slide, code, x=0.6, y=2.1, w=SW-1.2, h=None, size=14):
    lines = code.split("\n")
    if h is None: h = 0.32*len(lines) + 0.4
    _rect(slide, x, y, w, h, fill=CODEBG, rounded=True, radius=0.05)
    _dot(slide, x+0.22, y+0.2, 0.12, RGBColor(0xF2,0x6D,0x5B))
    _dot(slide, x+0.42, y+0.2, 0.12, RGBColor(0xF4,0xBF,0x4F))
    _dot(slide, x+0.62, y+0.2, 0.12, RGBColor(0x62,0xC5,0x54))
    runs = [[(ln if ln else " ", {"font": MONO, "size": size, "color": CODEFG})] for ln in lines]
    _txt(slide, x+0.3, y+0.5, w-0.6, h-0.6, runs, space_after=2, line_spacing=1.06, wrap=False)


def image_center(slide, path, x=0.6, y=2.0, w=SW-1.2, h=SH-2.6, caption=None, block=0):
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw/ih; box_ar = w/h
    if ar > box_ar: dw = w; dh = w/ar
    else: dh = h; dw = h*ar
    dx = x + (w-dw)/2; dy = y + (h-dh)/2
    slide.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    if caption:
        _txt(slide, x, y+h+0.02, w, 0.3, [[(caption, {"size": 12, "italic": True, "color": MUTED})]],
             align=PP_ALIGN.CENTER)


def notes(slide, script):
    slide.notes_slide.notes_text_frame.text = script


# ---------- PDF COMPANHEIRO ----------
def render_companion_pdf(pptx_path, slides_meta, out_pdf, workdir):
    import weasyprint
    os.makedirs(workdir, exist_ok=True)
    for f in glob.glob(os.path.join(workdir, "slide-*.png")): os.remove(f)
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", workdir, pptx_path],
                   check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=300)
    pdf = os.path.join(workdir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    subprocess.run(["pdftoppm", "-png", "-r", "110", pdf, os.path.join(workdir, "slide")], check=True, timeout=300)
    pngs = sorted(glob.glob(os.path.join(workdir, "slide-*.png")))
    def esc(s): return html.escape(s).replace("\n", "<br/>")
    FAMHEX = {k: ("#%02x%02x%02x" % (v[0][0], v[0][1], v[0][2])) for k, v in FAM.items()}
    def colhex(block): return FAMHEX[BLOCK_FAM.get(block, "indigo")]
    css = """
    @page { size: A4; margin: 1.2cm 1.4cm; }
    body { font-family:'Liberation Sans','DejaVu Sans',sans-serif; color:#161e33; font-size:11px; }
    h1 { font-size:19px; margin:0 0 2px 0; }
    .card { page-break-inside: avoid; margin-bottom: 15px; padding-bottom:8px; border-bottom:2px solid #eef; }
    .slideimg { width:100%; border-radius:8px; box-shadow:0 2px 8px rgba(20,26,50,.16); }
    .meta { display:flex; justify-content:space-between; align-items:center; margin:6px 0 4px 0; }
    .badge { color:#fff; padding:2px 10px; border-radius:12px; font-weight:bold; font-size:12px; }
    .tempo { padding:2px 10px; border-radius:12px; font-weight:bold; font-size:11px; background:#eef; color:#333; }
    .script { padding:9px 13px; border-radius:0 8px 8px 0; font-size:12px; line-height:1.55; background:#f7f8fc; }
    .script b { }
    .cover { text-align:center; padding-top:40px; }
    .cover h1 { font-size:27px; color:#4f46e5; } .cover p { color:#5a667e; }
    """
    body = ['<div class="cover"><h1>Roteiro do Apresentador</h1>'
            '<p>Engenharia de IA, Agentes e Desenvolvimento Orientado a Especificação — 120 minutos, 68 slides</p>'
            '<p>Para cada slide: a imagem, o tempo previsto e o texto a ser falado.</p></div>'
            '<div style="page-break-after:always"></div>']
    # pareamento POSICIONAL: o i-ésimo png corresponde ao i-ésimo item de slides_meta
    for i, png in enumerate(pngs, start=1):
        m = slides_meta[i-1] if i-1 < len(slides_meta) else {}
        col = colhex(m.get("block", 0))
        label = m.get("label") or ("S%d" % m.get("n", i))
        title = m.get("title", "Slide %d" % i)
        tempo = ("%s min" % _fmtmin(m["minutes"])) if m.get("minutes") is not None else "—"
        acum = ("  ·  acumulado %s min" % _fmtmin(m["acum"])) if m.get("acum") is not None else ""
        script = m.get("script", "")
        abs_png = "file://" + os.path.abspath(png)
        body.append(
            '<div class="card">'
            '<div class="meta"><span class="badge" style="background:%s">%s</span>'
            '<span class="tempo">%s%s</span></div>'
            '<h1 style="color:%s">%s</h1>'
            '<img class="slideimg" src="%s"/>'
            '<div class="meta"><span style="color:%s;font-weight:bold">▸ Roteiro de fala</span></div>'
            '<div class="script" style="border-left:4px solid %s">%s</div></div>'
            % (col, label, tempo, acum, col, esc(title), abs_png, col, col, script))
    doc = "<html><head><meta charset='utf-8'><style>%s</style></head><body>%s</body></html>" % (css, "".join(body))
    weasyprint.HTML(string=doc, base_url=workdir).write_pdf(out_pdf)
    return out_pdf, len(pngs)


def _fmtmin(x):
    if x is None: return "—"
    if abs(x - round(x)) < 1e-6: return str(int(round(x)))
    return ("%.1f" % x).replace(".", ",")
